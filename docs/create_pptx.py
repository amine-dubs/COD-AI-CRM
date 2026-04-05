"""
Generate PowerPoint presentation for COD-CRM Project.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DIAGRAMS_DIR = os.path.join(os.path.dirname(__file__), "diagrams")
OUTPUT = os.path.join(os.path.dirname(__file__), "presentation.pptx")

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF4)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
MED_GRAY = RGBColor(0x7F, 0x8C, 0x8D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left, top, width, font_size=16, color=DARK_GRAY, spacing=Pt(6)):
    """Add bulleted list to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_bar(slide, top=Inches(1.6)):
    """Add thin colored accent bar."""
    add_shape(slide, Inches(0.7), top, Inches(1.5), Pt(4), ACCENT_BLUE)


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BLUE)

# Accent rectangle at top
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Conception et Developpement d'un", font_size=28, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
             "CRM Intelligent", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "Integrant un Module de Prediction des Ventes\nBase sur le Machine Learning pour le E-Commerce",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Separator line
add_shape(slide, Inches(4.5), Inches(4.6), Inches(4.3), Pt(2), ACCENT_BLUE)

# Enterprise
add_text_box(slide, Inches(1), Inches(4.9), Inches(11), Inches(0.4),
             "Entreprise : BAT PROJET ENGINEERING  |  Represente par Mr. ATALLAH Mohammed",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Team and supervisor
add_text_box(slide, Inches(0.5), Inches(5.5), Inches(6), Inches(1.2),
             "Rezaiguia Soltane Tadj Eddine (DS)    Bellatreche Mohamed Amine (DS)\n"
             "Khelifi Ayyoub (CS)    Brahim Soheib (AI)",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7.5), Inches(5.5), Inches(5), Inches(0.8),
             "Encadre par : Dr. Khellat Souad",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Projet de Fin d'Etudes  |  Specialite : Data Science  |  Annee 2024-2025",
             font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

TOTAL_SLIDES = 16

# ════════════════════════════════════════════════════════════
# SLIDE 2: Plan / Agenda
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(6), Inches(0.8),
             "Plan de la Presentation", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

chapters = [
    ("01", "Contexte et Problematique", "E-commerce COD en Algerie, defis, objectifs"),
    ("02", "Architecture du Systeme", "Architecture globale, base de donnees, microservices"),
    ("03", "Integration IA", "Flux IA, pipeline ML, cycle de vie commande"),
    ("04", "Modeles de Machine Learning", "Risk scoring, segmentation, prevision demande"),
    ("05", "Resultats et Evaluation", "Comparaison des modeles, metriques, interpretation"),
    ("06", "Demonstration API", "Endpoints, exemples de requetes/reponses"),
]

for i, (num, title, desc) in enumerate(chapters):
    y = Inches(2.1) + Inches(0.85) * i
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT_BLUE
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, Inches(1.8), y - Pt(2), Inches(5), Inches(0.35),
                 title, font_size=18, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(1.8), y + Pt(20), Inches(5), Inches(0.3),
                 desc, font_size=12, color=MED_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 3: Contexte et Problematique
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.8),
             "Contexte et Problematique", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

# Left column - COD Model
box1 = add_rounded_shape(slide, Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.8), LIGHT_GRAY, ACCENT_BLUE)
box1.line.width = Pt(1)

add_text_box(slide, Inches(1.0), Inches(2.1), Inches(5), Inches(0.5),
             "Le modele COD en Algerie", font_size=20, bold=True, color=ACCENT_BLUE)

items_left = [
    "90%+ des transactions en Cash-on-Delivery",
    "58 wilayas, 3 zones d'expedition",
    "Monnaie : Dinar Algerien (DZD)",
    "Population jeune et connectee",
]
add_bullet_slide_content(slide, items_left, Inches(1.0), Inches(2.7), Inches(5.2), font_size=15)

# Right column - Problems
box2 = add_rounded_shape(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.8), LIGHT_GRAY, ACCENT_RED)
box2.line.width = Pt(1)

add_text_box(slide, Inches(7.1), Inches(2.1), Inches(5), Inches(0.5),
             "Problematique", font_size=20, bold=True, color=ACCENT_RED)

items_right = [
    "30-50% des commandes echouent (retour, refus, absence)",
    "Couts de transport aller-retour sans recette",
    "Confirmations telephoniques manuelles",
    "Aucune capacite predictive dans les CRM actuels",
]
add_bullet_slide_content(slide, items_right, Inches(7.1), Inches(2.7), Inches(5.2), font_size=15)

# Stat banner at bottom
stat_bar = add_shape(slide, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.8), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.6),
             "Objectif : Integrer le Machine Learning dans le CRM pour transformer les donnees en decisions actionnables",
             font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 4: Architecture Globale (diagram)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.8),
             "Architecture Globale du Systeme", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

img_path = os.path.join(DIAGRAMS_DIR, "architecture_globale.drawio.png")
slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.9), Inches(10.3), Inches(5.2))

add_slide_number(slide, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 5: Schema Base de Donnees (diagram)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.8),
             "Schema Relationnel de la Base de Donnees", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

add_text_box(slide, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
             "12 tables  |  Architecture multi-tenant (store_id)  |  MySQL 8.0  |  UTF8MB4",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.LEFT)

img_path = os.path.join(DIAGRAMS_DIR, "schema_base_donnees.drawio.png")
slide.shapes.add_picture(img_path, Inches(0.5), Inches(2.2), Inches(12.3), Inches(5.0))

add_slide_number(slide, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 6: AI Integration (diagram)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_PURPLE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.8),
             "Integration IA dans le CRM", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide, top=Inches(1.6))

img_path = os.path.join(DIAGRAMS_DIR, "ai_integration.drawio.png")
slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.3))

add_slide_number(slide, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 7: ML Pipeline (diagram)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_GREEN)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.8),
             "Pipeline Machine Learning", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

img_path = os.path.join(DIAGRAMS_DIR, "ml_pipeline.drawio.png")
slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.3))

add_slide_number(slide, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 8: Cycle de vie commande (diagram)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_ORANGE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.8),
             "Cycle de Vie d'une Commande COD", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

img_path = os.path.join(DIAGRAMS_DIR, "cycle_vie_commande.drawio.png")
slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.3))

add_slide_number(slide, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 9: Risk Prediction Model
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.8),
             "Prediction du Risque de Livraison", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

# Problem formulation
add_text_box(slide, Inches(0.7), Inches(2.0), Inches(11), Inches(0.4),
             "Classification binaire : predire si une commande sera livree (y=1) ou echouera (y=0)",
             font_size=16, color=DARK_GRAY)

# Feature categories
feat_data = [
    ("Temporelles", "6 features", "Heure, jour, mois, jours feries/weekend, trimestre", ACCENT_BLUE),
    ("Valeur", "4 features", "Montant, sous-total, frais, ratio", ACCENT_GREEN),
    ("Client", "3 features", "Recurrence, nb commandes, panier moyen", ACCENT_ORANGE),
    ("Paiement", "6 features", "Boleto/COD, CB, debit, voucher, mensualites", ACCENT_PURPLE),
    ("Qualite/Geo", "12 features", "Photos, desc., volume, poids, region, vendeur", ACCENT_RED),
]

for i, (cat, count, desc, color) in enumerate(feat_data):
    x = Inches(0.7) + Inches(2.45) * i
    box = add_rounded_shape(slide, x, Inches(2.8), Inches(2.3), Inches(1.5), LIGHT_GRAY, color)
    box.line.width = Pt(2)
    add_text_box(slide, x + Pt(8), Inches(2.9), Inches(2.1), Inches(0.35),
                 cat, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Pt(8), Inches(3.3), Inches(2.1), Inches(0.25),
                 count, font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Pt(8), Inches(3.65), Inches(2.1), Inches(0.5),
                 desc, font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Ensemble approach
add_text_box(slide, Inches(0.7), Inches(4.6), Inches(11), Inches(0.4),
             "Ensemble par vote pondere (Soft Voting) de 3 modeles :", font_size=16, bold=True, color=DARK_BLUE)

models_info = [
    ("CatBoost", "Gestion native des\nvariables categorielles"),
    ("LightGBM", "GOSS + EFB\nCroissance leaf-wise"),
    ("XGBoost", "Regularisation L1/L2\nOptimisation Newton"),
]

for i, (name, desc) in enumerate(models_info):
    x = Inches(0.7) + Inches(4.1) * i
    box = add_rounded_shape(slide, x, Inches(5.2), Inches(3.8), Inches(1.5), DARK_BLUE)
    add_text_box(slide, x + Pt(8), Inches(5.3), Inches(3.6), Inches(0.4),
                 name, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Pt(8), Inches(5.8), Inches(3.6), Inches(0.7),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 10: Risk Model RESULTS - Comparison Table
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.8),
             "Resultats : Comparaison des Modeles de Classification", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

add_text_box(slide, Inches(0.7), Inches(1.8), Inches(11), Inches(0.4),
             "Test : 19 543 echantillons (80/20 split stratifie)  |  31 features  |  ADASYN + Optuna", font_size=14, color=MED_GRAY)

# Table header
headers = ["Modele", "AUC-ROC", "Accuracy", "Precision", "Rappel", "F1-Score"]
col_widths = [2.2, 1.7, 1.7, 1.7, 1.7, 1.7]
x_start = Inches(0.9)

for j, (h, w) in enumerate(zip(headers, col_widths)):
    x = x_start + Inches(sum(col_widths[:j]))
    box = add_shape(slide, x, Inches(2.3), Inches(w), Inches(0.5), DARK_BLUE)
    add_text_box(slide, x, Inches(2.33), Inches(w), Inches(0.45),
                 h, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Table rows
rows = [
    ("CatBoost", "0.9957", "0.9997", "0.9997", "1.0000", "0.9999"),
    ("LightGBM", "0.9974", "0.9997", "0.9997", "1.0000", "0.9999"),
    ("XGBoost", "0.9967", "0.9996", "0.9997", "0.9999", "0.9998"),
    ("Ensemble", "0.9961", "0.9997", "0.9997", "1.0000", "0.9999"),
]

for i, row in enumerate(rows):
    y = Inches(2.8) + Inches(0.5) * i
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    is_last = i == len(rows) - 1
    for j, (val, w) in enumerate(zip(row, col_widths)):
        x = x_start + Inches(sum(col_widths[:j]))
        fill = ACCENT_BLUE if is_last else bg
        box = add_shape(slide, x, y, Inches(w), Inches(0.5), fill)
        if is_last:
            box.line.color.rgb = ACCENT_BLUE
        txt_color = WHITE if is_last else DARK_GRAY
        bld = is_last
        add_text_box(slide, x, y + Pt(3), Inches(w), Inches(0.45),
                     val, font_size=14, bold=bld, color=txt_color, alignment=PP_ALIGN.CENTER)

# Confusion matrix
add_text_box(slide, Inches(0.9), Inches(5.2), Inches(5), Inches(0.4),
             "Matrice de Confusion (Ensemble)", font_size=18, bold=True, color=DARK_BLUE)

# Simple confusion matrix
cm_data = [
    ("", "Predit: Echouee", "Predit: Livree"),
    ("Reel: Echouee", "242 (VN)", "5 (FP)"),
    ("Reel: Livree", "0 (FN)", "19 296 (VP)"),
]

for i, row in enumerate(cm_data):
    for j, val in enumerate(row):
        x = Inches(0.9) + Inches(2.0) * j
        y = Inches(5.7) + Inches(0.45) * i
        if i == 0 or j == 0:
            bg_c = DARK_BLUE if (i == 0 and j > 0) or (j == 0 and i > 0) else MED_GRAY
            txt_c = WHITE
            bld = True
        else:
            if (i == 1 and j == 1) or (i == 2 and j == 2):
                bg_c = RGBColor(0xD5, 0xE8, 0xD4)  # green for correct
            else:
                bg_c = RGBColor(0xF8, 0xCE, 0xCC)  # red for errors
            txt_c = DARK_GRAY
            bld = False

        if i == 0 and j == 0:
            bg_c = WHITE
            txt_c = WHITE

        box = add_shape(slide, x, y, Inches(2.0), Inches(0.45), bg_c)
        add_text_box(slide, x, y + Pt(2), Inches(2.0), Inches(0.4),
                     val, font_size=12, bold=bld, color=txt_c, alignment=PP_ALIGN.CENTER)

# Key insight
insight_box = add_rounded_shape(slide, Inches(7.0), Inches(5.3), Inches(5.5), Inches(1.5), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE)
insight_box.line.width = Pt(2)
add_text_box(slide, Inches(7.2), Inches(5.4), Inches(5.1), Inches(0.3),
             "Points cles", font_size=14, bold=True, color=ACCENT_ORANGE)
items_insight = [
    "AUC-ROC : 0.9961 (31 features, ADASYN, Optuna)",
    "Rappel echecs : 98% | Precision : 100% | F1 = 0.99",
    "Seuil ensemble : 0.9805 (Youden's J statistic)",
    "Seuils individuels : F1-max par modele (precision_recall_curve)",
]
add_bullet_slide_content(slide, items_insight, Inches(7.2), Inches(5.8), Inches(5.1), font_size=12, color=DARK_GRAY, spacing=Pt(2))

add_slide_number(slide, 10, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 11: Segmentation Client
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_PURPLE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.8),
             "Segmentation Client (Hybride HDBSCAN / KMeans)", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

# RFM explanation
add_text_box(slide, Inches(0.7), Inches(1.8), Inches(11), Inches(0.4),
             "Analyse RFM + Clustering adaptatif : HDBSCAN (>=1000 clients) ou KMeans avec auto-K (<1000 clients)",
             font_size=15, color=MED_GRAY)

# Segment cards
segments = [
    ("VIP", "3 673", "3.8%", "Haute valeur, achats frequents", ACCENT_GREEN),
    ("Loyal", "252", "0.3%", "Clients reguliers", ACCENT_BLUE),
    ("A Risque", "2 736", "2.8%", "Declin d'activite", ACCENT_ORANGE),
    ("Perdu", "441", "0.5%", "Inactifs", ACCENT_RED),
    ("Standard", "88 994", "92.6%", "Clients occasionnels", MED_GRAY),
]

for i, (name, count, pct, desc, color) in enumerate(segments):
    x = Inches(0.5) + Inches(2.5) * i
    # Card
    card = add_rounded_shape(slide, x, Inches(2.5), Inches(2.3), Inches(2.8), WHITE, color)
    card.line.width = Pt(2)
    card.shadow.inherit = False

    # Color header bar
    add_shape(slide, x, Inches(2.5), Inches(2.3), Inches(0.5), color)
    add_text_box(slide, x, Inches(2.52), Inches(2.3), Inches(0.45),
                 name, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Count
    add_text_box(slide, x, Inches(3.2), Inches(2.3), Inches(0.5),
                 count, font_size=28, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.7), Inches(2.3), Inches(0.3),
                 f"({pct})", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(4.2), Inches(2.3), Inches(0.6),
                 desc, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Total banner
total_bar = add_shape(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5), DARK_BLUE)
add_text_box(slide, Inches(0.5), Inches(5.63), Inches(12.3), Inches(0.45),
             "Total : 96 096 clients (Olist = HDBSCAN) | Petites BDD = KMeans auto-K (silhouette optimisee)",
             font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Actions
add_text_box(slide, Inches(0.7), Inches(6.3), Inches(12), Inches(0.4),
             "Actions : VIP -> Fidelite & offres exclusives  |  A Risque -> Campagnes de reactivation  |  Standard -> Strategies de conversion",
             font_size=13, color=DARK_GRAY)

add_slide_number(slide, 11, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 12: Demand Forecasting
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_GREEN)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.8),
             "Prevision de la Demande (LightGBM + Calendrier Islamique)", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

add_text_box(slide, Inches(0.7), Inches(1.8), Inches(11), Inches(0.4),
             "Benchmark de 5 modeles covariants  |  Serie temporelle : CA quotidien DZD sur 714 jours  |  Evaluation out-of-sample",
             font_size=14, color=MED_GRAY)

# Benchmark table header
bench_headers = [("Modele", 3.0), ("MAE (DZD)", 2.5), ("RMSE (DZD)", 2.5), ("MAPE", 2.0), ("vs MA7", 2.0)]
for j, (h, w) in enumerate(bench_headers):
    x = Inches(0.6) + Inches(sum([bw for _, bw in bench_headers[:j]]))
    box = add_shape(slide, x, Inches(2.3), Inches(w), Inches(0.45), DARK_BLUE)
    add_text_box(slide, x, Inches(2.32), Inches(w), Inches(0.4),
                 h, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Benchmark rows (30-day horizon, recursive forecasting)
bench_rows = [
    ("Moyenne Mobile (7j)", "363 075", "466 570", "163.9%", "baseline", LIGHT_GRAY),
    ("Prophet + Fetes islamiques", "376 292", "447 736", "150.6%", "-3.6%", WHITE),
    ("Chronos-T5-Small", "338 231", "395 363", "132.9%", "+6.8%", LIGHT_GRAY),
    ("LightGBM + lags + holidays", "318 741", "416 501", "148.0%", "+12.2%", ACCENT_GREEN),
]

for i, (name, mae, rmse, mape, imp, bg) in enumerate(bench_rows):
    y = Inches(2.75) + Inches(0.42) * i
    is_best = (bg == ACCENT_GREEN)
    for j, (val, w) in enumerate(zip([name, mae, rmse, mape, imp], [bw for _, bw in bench_headers])):
        x = Inches(0.6) + Inches(sum([bw for _, bw in bench_headers[:j]]))
        fill = ACCENT_GREEN if is_best else bg
        box = add_shape(slide, x, y, Inches(w), Inches(0.42), fill)
        txt_color = WHITE if is_best else DARK_GRAY
        bld = is_best
        add_text_box(slide, x, y + Pt(2), Inches(w), Inches(0.38),
                     val, font_size=12, bold=bld, color=txt_color, alignment=PP_ALIGN.CENTER)

# LightGBM features (left)
add_text_box(slide, Inches(0.7), Inches(4.6), Inches(5.5), Inches(0.4),
             "LightGBM + Calendrier Algerien", font_size=18, bold=True, color=DARK_BLUE)

config_items = [
    "20 features : lags(1,7,14,28), rolling stats, calendrier",
    "Islamiques : Ramadan, Aid el-Fitr, Aid el-Adha, Mawlid",
    "Nationaux : 1er jan, Yennayer, 1er mai, 5 juil, 1er nov",
    "is_weekend = ven-sam + jours feries islamiques + nationaux",
    "Prevision recursive multi-step (pas de fuite)",
    "4 series : global + 3 categories produit",
]
add_bullet_slide_content(slide, config_items, Inches(0.7), Inches(5.05), Inches(5.5), font_size=12, spacing=Pt(2))

# Key insight (right)
insight_box2 = add_rounded_shape(slide, Inches(6.8), Inches(4.55), Inches(5.7), Inches(2.6), RGBColor(0xD5, 0xE8, 0xD4), ACCENT_GREEN)
insight_box2.line.width = Pt(2)
add_text_box(slide, Inches(7.0), Inches(4.6), Inches(5.3), Inches(0.3),
             "Pourquoi LightGBM ?", font_size=14, bold=True, color=ACCENT_GREEN)
insight_items = [
    "Meilleur MAE : 318 741 DZD (+12.2% vs baseline)",
    "Supporte les covariables (calendrier algerien)",
    "Modeles fondamentaux (Chronos, TimesFM, MOIRAI)",
    "= zero-shot sans covariables -> inadaptes",
    "714 jours favorisent LightGBM vs deep learning",
]
add_bullet_slide_content(slide, insight_items, Inches(7.0), Inches(4.95), Inches(5.3), font_size=12, color=DARK_GRAY, spacing=Pt(2))

add_slide_number(slide, 12, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 13: API Demo
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.8),
             "API REST : Endpoints du Module ML", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

add_text_box(slide, Inches(0.7), Inches(1.8), Inches(11), Inches(0.4),
             "FastAPI (Python) - 15 endpoints - Port 8001  |  PHP Backend - Port 8000", font_size=15, color=MED_GRAY)

# Endpoint categories with icons
categories = [
    ("Prediction", ACCENT_RED, [
        "POST /api/predict/order-risk",
        "POST /api/predict/order-risk/batch",
        "GET  /api/predict/model-info",
    ]),
    ("Segmentation", ACCENT_PURPLE, [
        "GET  /api/segment/customers",
        "GET  /api/segment/summary",
    ]),
    ("Prevision", ACCENT_GREEN, [
        "GET  /api/forecast/demand",
        "GET  /api/forecast/categories",
    ]),
    ("Reentrainement", ACCENT_ORANGE, [
        "POST /api/v1/ai/retrain (depuis BDD)",
        "POST /api/retrain/upload-and-train",
        "POST /api/retrain/restore-defaults",
    ]),
]

for i, (cat_name, color, endpoints) in enumerate(categories):
    y = Inches(2.3) + Inches(1.2) * i

    # Category label
    cat_box = add_rounded_shape(slide, Inches(0.7), y, Inches(2.2), Inches(0.45), color)
    add_text_box(slide, Inches(0.7), y + Pt(2), Inches(2.2), Inches(0.4),
                 cat_name, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Endpoints
    for j, ep in enumerate(endpoints):
        add_text_box(slide, Inches(3.2), y + Inches(0.3) * j, Inches(6), Inches(0.3),
                     ep, font_size=13, color=DARK_GRAY, font_name="Consolas")

# Example response box
resp_box = add_rounded_shape(slide, Inches(8.5), Inches(2.3), Inches(4.3), Inches(4.5), RGBColor(0x2C, 0x3E, 0x50))
add_text_box(slide, Inches(8.7), Inches(2.4), Inches(3.9), Inches(0.3),
             "Exemple de reponse", font_size=12, bold=True, color=ACCENT_BLUE)

resp_text = '{\n  "score": 67.9,\n  "category": "medium",\n  "success_probability": 0.68,\n  "reasons": [\n    "Region: low success",\n    "Category: high return"\n  ],\n  "recommendation":\n    "Monitor closely",\n  "model_scores": {\n    "catboost": 92.1,\n    "lightgbm": 12.7,\n    "xgboost": 99.0\n  }\n}'
add_text_box(slide, Inches(8.7), Inches(2.8), Inches(3.9), Inches(3.8),
             resp_text, font_size=10, color=ACCENT_GREEN, font_name="Consolas")

add_slide_number(slide, 13, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 14: Retraining & Production Readiness
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_GREEN)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.8),
             "Reentrainement et Mise en Production", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

# Left: Retraining flow
left_box = add_rounded_shape(slide, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0), LIGHT_GRAY, ACCENT_GREEN)
left_box.line.width = Pt(2)

add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(0.4),
             "Pipeline de reentrainement", font_size=18, bold=True, color=ACCENT_GREEN)

retrain_steps = [
    "1. Clic sur 'Reentrainer depuis la base de donnees'",
    "2. Extraction automatique des commandes (PHP -> CSV)",
    "3. Backup automatique des modeles actuels",
    "4. Optimisation Optuna (40 essais bayesiens)",
    "5. Entrainement des 3 modeles (params optimises)",
    "6. Sauvegarde metriques + parametres optimaux",
    "7. Rechargement automatique dans le service",
    "8. Rollback possible via restore-defaults",
]
add_bullet_slide_content(slide, retrain_steps, Inches(0.7), Inches(2.6), Inches(5.6), font_size=13, spacing=Pt(3))

# Right: What's saved
right_box = add_rounded_shape(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(2.3), LIGHT_GRAY, ACCENT_BLUE)
right_box.line.width = Pt(2)

add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.4),
             "Modeles et metriques sauvegardes", font_size=18, bold=True, color=ACCENT_BLUE)

saved_items = [
    "risk_ensemble.joblib (7.6 MB) - Ensemble CatBoost+LightGBM+XGBoost",
    "segmenter.joblib (9.7 MB) - HDBSCAN ou KMeans (adaptatif)",
    "forecaster_lgbm.joblib (1.8 MB) - LightGBM + calendrier islamique",
    "metrics.json (4 KB) - Toutes les metriques d'evaluation",
]
add_bullet_slide_content(slide, saved_items, Inches(7.0), Inches(2.6), Inches(5.6), font_size=12, spacing=Pt(2))

# Bottom right: API format
format_box = add_rounded_shape(slide, Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.5), RGBColor(0x2C, 0x3E, 0x50))

add_text_box(slide, Inches(7.0), Inches(4.6), Inches(5.6), Inches(0.3),
             "2 methodes de reentrainement", font_size=13, bold=True, color=ACCENT_BLUE)

format_items = [
    "Methode 1 (recommandee) : depuis la base de donnees",
    "  -> POST /api/v1/ai/retrain (un clic, automatique)",
    "Methode 2 : upload CSV personnalise",
    "  -> POST /api/retrain/upload-and-train",
    "Optuna auto-optimise les hyperparametres",
]
add_bullet_slide_content(slide, format_items, Inches(7.0), Inches(5.0), Inches(5.6), font_size=11, color=LIGHT_GRAY, spacing=Pt(1))

add_slide_number(slide, 14, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 15: Security & Workflow Actions
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_RED)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.8),
             "Securite et Workflow Automatise", font_size=34, bold=True, color=DARK_BLUE)
add_accent_bar(slide)

# Left: Security measures
left_box = add_rounded_shape(slide, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.8), LIGHT_GRAY, ACCENT_RED)
left_box.line.width = Pt(2)

add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(0.4),
             "Mesures de securite", font_size=18, bold=True, color=ACCENT_RED)

security_items = [
    "Authentification : JWT + RBAC (6 roles) + cle API ML",
    "Injection SQL : Requetes preparees (PDO parametrisees)",
    "CSRF : Header X-Requested-With obligatoire",
    "Rate Limiting : 5 req/min (auth), 60/min (general)",
    "Upload : Limite 50 Mo, validation MIME, anti-path traversal",
    "Debug : APP_DEBUG=false en production",
    "Multi-tenant : Isolation stricte par store_id",
]
add_bullet_slide_content(slide, security_items, Inches(0.7), Inches(2.6), Inches(5.6), font_size=13, spacing=Pt(3))

# Right top: Workflow actions
right_box = add_rounded_shape(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(2.5), LIGHT_GRAY, ACCENT_GREEN)
right_box.line.width = Pt(2)

add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.4),
             "Actions de workflow automatisees", font_size=18, bold=True, color=ACCENT_GREEN)

# Workflow table header
workflow_headers = ["Score", "Action", "Description"]
for i, h in enumerate(workflow_headers):
    w = Inches(1.2) if i == 0 else Inches(1.5) if i == 1 else Inches(2.9)
    x = Inches(7.0) + sum([Inches(1.2), Inches(1.5), Inches(0)][:i])
    add_shape(slide, x, Inches(2.55), w, Inches(0.35), DARK_BLUE)
    add_text_box(slide, x, Inches(2.57), w, Inches(0.3),
                 h, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Workflow table rows
workflow_rows = [
    ("70-100%", "Auto-Approve", "Risque faible - traitement auto", ACCENT_GREEN),
    ("30-70%", "Revue manuelle", "Verification humaine requise", ACCENT_ORANGE),
    ("0-30%", "Signalement", "Rejet recommande", ACCENT_RED),
]
for row_idx, (score, action, desc, color) in enumerate(workflow_rows):
    y = Inches(2.95) + Inches(0.4) * row_idx
    add_shape(slide, Inches(7.0), y, Inches(1.2), Inches(0.35), WHITE, color)
    add_text_box(slide, Inches(7.0), y + Pt(2), Inches(1.2), Inches(0.3),
                 score, font_size=10, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.2), y + Pt(2), Inches(1.5), Inches(0.3),
                 action, font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.7), y + Pt(2), Inches(2.9), Inches(0.3),
                 desc, font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Right bottom: Youden explanation
youden_box = add_rounded_shape(slide, Inches(6.8), Inches(4.7), Inches(6.0), Inches(2.1), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE)
youden_box.line.width = Pt(2)

add_text_box(slide, Inches(7.0), Inches(4.8), Inches(5.6), Inches(0.4),
             "Seuil optimal : Youden's J statistic", font_size=14, bold=True, color=ACCENT_ORANGE)

youden_items = [
    "J = Sensibilite + Specificite - 1",
    "Maximise simultanement vrais positifs et negatifs",
    "Adapte au desequilibre de classes (2-3% echecs)",
    "Seuil calcule : 0.9805 (cours ROC)",
]
add_bullet_slide_content(slide, youden_items, Inches(7.0), Inches(5.2), Inches(5.6), font_size=11, color=DARK_GRAY, spacing=Pt(1))

add_slide_number(slide, 15, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SLIDE 16: Conclusion & Perspectives
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), ACCENT_BLUE)

add_text_box(slide, Inches(0.7), Inches(0.5), Inches(10), Inches(0.8),
             "Conclusion et Perspectives", font_size=34, bold=True, color=WHITE)
add_shape(slide, Inches(0.7), Inches(1.3), Inches(1.5), Pt(4), ACCENT_BLUE)

# Contributions
add_text_box(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.4),
             "Contributions", font_size=22, bold=True, color=ACCENT_BLUE)

contributions = [
    "CRM complet : 12 tables, 11 modules, multi-tenant",
    "Ensemble ML (AUC=0.9961) : seuil Youden optimal, rappel echecs 98%",
    "Workflow automatise : auto-approve / revue manuelle / signalement",
    "Segmentation hybride HDBSCAN/KMeans adaptative selon taille BDD",
    "Prevision LightGBM + Calendrier algerien + date de depart flexible",
    "Securite : JWT+RBAC, rate limiting, CSRF, requetes preparees",
    "Interface multilingue (AR/FR/EN + RTL) + LLM Gemini",
]
add_bullet_slide_content(slide, contributions, Inches(0.7), Inches(2.2), Inches(5.8), font_size=14, color=LIGHT_GRAY, spacing=Pt(3))

# Perspectives
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4),
             "Perspectives", font_size=22, bold=True, color=ACCENT_ORANGE)

perspectives = [
    "Deploiement sur donnees reelles algeriennes",
    "Evenements commerciaux (rentree, promotions)",
    "Apprentissage en ligne (mise a jour continue)",
    "A/B testing de l'impact des recommandations ML",
]
add_bullet_slide_content(slide, perspectives, Inches(7.0), Inches(2.2), Inches(5.8), font_size=14, color=LIGHT_GRAY, spacing=Pt(3))

# Thank you
add_shape(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(1.5), ACCENT_BLUE)
add_text_box(slide, Inches(3.5), Inches(5.7), Inches(6.3), Inches(0.6),
             "Merci pour votre attention", font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.4),
             "Questions ?", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 16, TOTAL_SLIDES)

# ════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {TOTAL_SLIDES}")
